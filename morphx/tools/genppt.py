# Standard libraries
import json
import logging
import re
from datetime import datetime
from pathlib import Path

# 3pps
import torch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Own modules
from morphx.inputs import PPT_INPUT_CONTENT
from morphx.models import SlideModel
from morphx.utils import (
    ContentParsingError,
    PPTGeneratorError,
    SlideGenerationError,
    create_logger,
    read_json_config,
    validate_slide_content,
)


class PPTGenerator:
    """
    Production-ready PowerPoint generator with LLM enhancement.
    
    Features:
    - Robust error handling and validation
    - Section-by-section processing for better quality
    - Optional LLM enhancement
    - Memory optimization for RTX 3060
    - Comprehensive logging
    """
    
    DEFAULT_COLORS = {
        'primary': RGBColor(0, 51, 102),
        'secondary': RGBColor(51, 51, 51),
        'accent': RGBColor(150, 150, 150)
    }
    
    def __init__(self, slides_model: SlideModel, logger: logging.Logger | None):
        """
        Initialize the PPT generator.
        
        Args:
            slides_model: Modelo utilizado para la creación de slides.
            enable_logging: Enable detailed logging.
            
        Raises:
            ModelLoadError: If model fails to load
        """

        # Set attributes
        self.logger = logger
        self.slides_model = slides_model

    def extract_slides_from_content(self, content_text: str) -> list[str]:
        """
        Extract individual slide sections from content.
        
        Args:
            content_text: Full content with slide sections
            
        Returns:
            List of slide section strings
            
        Raises:
            ContentParsingError: If no slides found
        """

        self.logger.info("Extracting slide sections...")
        
        if not content_text or not content_text.strip():
            raise ContentParsingError("Content is empty")
        
        # Primary pattern: ### **Slide X** or ### Slide X
        slide_pattern = r'###\s*\*?Slide\s+\d+\*?\s*\n(.*?)(?=###\s*\*?Slide\s+\d+|$)'
        slides = re.findall(slide_pattern, content_text, re.DOTALL | re.IGNORECASE)
        
        # Fallback: split by ### headers
        if not slides:
            self.logger.info("Primary pattern failed, trying fallback...",)
            slides = re.split(r'###\s+', content_text)
            slides = [s.strip() for s in slides if s.strip() and len(s) > 50]
        
        if not slides:
            raise ContentParsingError("No valid slide sections found in content")
        
        self.logger.info(f"Found {len(slides)} slide sections")
        return slides

    def parse_slide_section(self, slide_text: str, slide_num: int) -> dict:
        """
        Parse a single slide section to extract structured information.
        
        Args:
            slide_text: Raw slide section text
            slide_num: Slide number for error reporting
            
        Returns:
            Dictionary with slide data
        """
        
        slide_data = {
            "identifier": f"SL{slide_num}",  # Default identifier
            "main_title": "",
            "subtitle": "",
            "key_points": []
        }
        
        try:
            # Extract identifier (case-insensitive)
            id_match = re.search(
                r'(?:Identifier|ID):\s*`?([^`\n]+)`?', 
                slide_text, 
                re.IGNORECASE
            )
            if id_match:
                slide_data["identifier"] = id_match.group(1).strip()
            
            # Extract main title
            title_match = re.search(
                r'(?:Main\s+Title|Title):\s*(.+?)(?:\n|$)', 
                slide_text, 
                re.IGNORECASE
            )
            if title_match:
                slide_data["main_title"] = title_match.group(1).strip()
            
            # Extract subtitle
            subtitle_match = re.search(
                r'Subtitle:\s*(.+?)(?:\n|$)', 
                slide_text, 
                re.IGNORECASE
            )
            if subtitle_match:
                slide_data["subtitle"] = subtitle_match.group(1).strip()
            
            # Extract key points section
            key_points_section = re.search(
                r'(?:Key\s+Points(?:\s+and\s+Description)?|Content):\s*\n(.*)', 
                slide_text, 
                re.DOTALL | re.IGNORECASE
            )
            
            if key_points_section:
                points_text = key_points_section.group(1)
                
                # Split by numbered points (handles various formats)
                point_pattern = r'(\d+)\.\s+\*?\*?([^\n*]+)\*?\*?\s*\n+(.*?)(?=\n\d+\.\s+|\Z)'
                points = re.findall(point_pattern, points_text, re.DOTALL)
                
                for _point_num, point_title, point_desc in points:
                    # Clean up description
                    desc_lines = []
                    for line in point_desc.strip().split('\n'):
                        line = line.strip()
                        # Remove leading asterisks and dashes
                        line = re.sub(r'^[\*\-]\s*', '', line)
                        if line and not line.startswith('---'):
                            desc_lines.append(line)
                    
                    description = ' '.join(desc_lines)
                    
                    if point_title.strip() or description:
                        slide_data["key_points"].append({
                            "point": point_title.strip(),
                            "description": description
                        })
            
            # Validation
            if not slide_data["subtitle"]:
                self.logger.info(f"Slide {slide_num}: Missing subtitle",)
            
            if not slide_data["key_points"]:
                self.logger.info(f"Slide {slide_num}: No key points found",)
            
        except Exception as e:
            self.logger.info(f"Error parsing slide {slide_num}: {str(e)}", "error")
        
        return slide_data

    def enhance_slide_with_llm(self, slide_data: dict) -> list[dict]:
        """
        Use LLM to enhance slide content.
        
        Args:
            slide_data: Parsed slide data
            
        Returns:
            Enhanced key points list
        """
        subtitle = slide_data.get('subtitle', 'Slide')
        self.logger.info(f"Enhancing: {subtitle[:60]}...")
        
        try:
            # Build context
            points_text = "\n".join([
                f"- {p['point']}: {p['description']}" 
                for p in slide_data['key_points']
            ])
            
            prompt = f"""
                You are an expert educational content creator. Refine this slide content to be more clear, academic, and engaging.

                Slide Topic: {subtitle}

                Current Content:
                {points_text}

                Instructions:
                1. Improve descriptions for clarity and educational value
                2. Use present tense and impersonal academic style
                3. Make explanations self-contained (no need for instructor)
                4. Keep each description concise (2-3 sentences max)
                5. Maintain all key concepts but enhance clarity

                Return ONLY a JSON array:
                [
                {{
                    "point": "Key Concept Title",
                    "description": "Enhanced academic description"
                }}
                ]
            """
            
            # Tokenize
            inputs = self.tokenizer(
                prompt,
                return_tensors="pt",
                padding=True,
                truncation=True,
                max_length=1536
            ).to(self.model.device)
            
            # Generate
            with torch.no_grad():
                output = self.model.generate(
                    **inputs,
                    do_sample=True,
                    temperature=0.4,
                    top_p=0.85,
                    repetition_penalty=1.1,
                    max_new_tokens=1024,
                    pad_token_id=self.tokenizer.pad_token_id,
                    eos_token_id=self.tokenizer.eos_token_id,
                )
            
            # Decode
            response = self.tokenizer.decode(
                output[0][inputs['input_ids'].shape[1]:], 
                skip_special_tokens=True
            )
            
            # Clear cache
            if torch.cuda.is_available():
                torch.cuda.empty_cache()
            
            # Parse JSON
            json_match = re.search(r'\[[\s\S]*?\]', response)
            if json_match:
                enhanced_points = json.loads(json_match.group())
                if isinstance(enhanced_points, list) and len(enhanced_points) > 0:
                    self.logger.info(f"  ✓ Enhanced {len(enhanced_points)} points")
                    return enhanced_points
            
        except json.JSONDecodeError as e:
            self.logger.info(f"  JSON parse error: {str(e)}",)
        except Exception as e:
            self.logger.info(f"  Enhancement error: {str(e)}",)
        
        # Return original on failure
        self.logger.info("  ⚠ Using original content",)
        return slide_data['key_points']

    def process_content_to_structure(
        self, 
        content_text: str, 
        use_llm_enhancement: bool = True
    ) -> dict:
        """
        Process content into structured slides.
        
        Args:
            content_text: Raw content text
            use_llm_enhancement: Whether to use LLM enhancement
            
        Returns:
            Structured presentation data
            
        Raises:
            ContentParsingError: If content cannot be parsed
        """
        
        self.logger.info("\n" + "="*60)
        self.logger.info("PROCESSING CONTENT")
        self.logger.info("="*60 + "\n")
        
        # Extract slides
        slide_sections = self.extract_slides_from_content(content_text)
        
        all_slides = []
        main_title = None
        
        for i, section in enumerate(slide_sections, 1):
            self.logger.info(f"\nProcessing slide {i}/{len(slide_sections)}...")
            
            # Parse section
            slide_data = self.parse_slide_section(section, i)
            
            # Store main title from first slide
            if i == 1 and slide_data['main_title']:
                main_title = slide_data['main_title']
            
            # Enhance with LLM if enabled
            if use_llm_enhancement and slide_data['key_points']:
                enhanced_points = self.enhance_slide_with_llm(slide_data)
                slide_data['key_points'] = enhanced_points
            
            # Only add slides with content
            if slide_data['key_points']:
                all_slides.append(slide_data)
            else:
                self.logger.info(f"  Skipping slide {i} (no content)",)
        
        if not all_slides:
            raise ContentParsingError("No valid slides generated")
        
        structure = {
            "main_title": main_title or "Presentation",
            "slides": all_slides
        }
        
        self.logger.info(f"\n{'='*60}")
        self.logger.info(f"COMPLETED: {len(all_slides)} slides processed")
        self.logger.info(f"{'='*60}\n")
        
        return structure

    def create_presentation(
        self, 
        structure: dict, 
        output_filename: str = "presentation.pptx"
    ) -> str:
        """
        Create PowerPoint presentation from structured content.
        
        Args:
            structure: Structured presentation data
            output_filename: Output file path
            
        Returns:
            Path to created presentation
            
        Raises:
            SlideGenerationError: If presentation creation fails
        """
        self.logger.info("Creating PowerPoint presentation...")
        
        try:
            # Validate output path
            output_path = Path(output_filename)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            main_title = structure.get("main_title", "Presentation")
            
            # Title slide
            self._create_title_slide(prs, main_title)
            
            # Content slides
            for slide_data in structure.get("slides", []):
                self._create_content_slide(prs, slide_data, main_title)
            
            # Save
            prs.save(str(output_path))
            self.logger.info(
                f"✓ Presentation saved: '{output_path}' ({len(structure['slides'])} slides)", 
                "success"
            )
            
            return str(output_path)
            
        except Exception as e:
            raise SlideGenerationError(f"Failed to create presentation: {str(e)}") from e

    def _create_title_slide(self, prs: Presentation, main_title: str):
        """Create title slide."""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = main_title
        subtitle.text = f"Academic Presentation\n{datetime.now().strftime('%B %Y')}"
        
        # Style
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.DEFAULT_COLORS['primary']

    def _create_content_slide(
        self, 
        prs: Presentation, 
        slide_data: dict, 
        main_title: str
    ):
        """Create a content slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Identifier (top right)
        id_box = slide.shapes.add_textbox(
            Inches(8.5), Inches(0.2), Inches(1.3), Inches(0.3)
        )
        id_frame = id_box.text_frame
        id_frame.text = slide_data.get("identifier", "")
        id_para = id_frame.paragraphs[0]
        id_para.alignment = PP_ALIGN.RIGHT
        id_para.font.size = Pt(10)
        id_para.font.color.rgb = self.DEFAULT_COLORS['accent']
        
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = main_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.DEFAULT_COLORS['primary']
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.1), Inches(9), Inches(0.7)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = slide_data.get("subtitle", "")
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(22)
        subtitle_para.font.color.rgb = self.DEFAULT_COLORS['secondary']
        subtitle_para.font.italic = True
        
        # Horizontal line
        line = slide.shapes.add_shape(
            1,  # Line shape
            Inches(0.5), Inches(1.9),
            Inches(9), Inches(0)
        )
        line.line.color.rgb = self.DEFAULT_COLORS['primary']
        line.line.width = Pt(2)
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.2), Inches(9), Inches(4.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for idx, point_data in enumerate(slide_data.get("key_points", [])):
            # Spacing
            if idx > 0:
                spacer = content_frame.add_paragraph()
                spacer.space_after = Pt(8)
            
            # Point title
            p = content_frame.add_paragraph()
            p.text = f"• {point_data.get('point', '')}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.DEFAULT_COLORS['primary']
            p.space_after = Pt(4)
            
            # Description
            desc = point_data.get('description', '')
            if desc and desc != point_data.get('point', ''):
                p_desc = content_frame.add_paragraph()
                p_desc.text = desc
                p_desc.font.size = Pt(13)
                p_desc.level = 1
                p_desc.space_after = Pt(8)
                p_desc.font.color.rgb = self.DEFAULT_COLORS['secondary']

    def generate_ppt(
        self, 
        content_text: str, 
        output_filename: str = "presentation.pptx", 
        use_llm_enhancement: bool = True
    ) -> str:
        """
        Complete pipeline: content -> structure -> PPT.
        
        Args:
            content_text: Input content with slide sections
            output_filename: Output PPT filename
            use_llm_enhancement: Whether to use LLM enhancement
            
        Returns:
            Path to created presentation
            
        Raises:
            PPTGeneratorError: If generation fails
        """
        try:
            structure = self.process_content_to_structure(
                content_text, 
                use_llm_enhancement
            )
            return self.create_presentation(structure, output_filename)
            
        except Exception as e:
            self.logger.info(f"Generation failed: {str(e)}", "error")
            raise

if __name__ == "__main__":
    try:
        # Primero leemos la configuracion de los modelos a partir de un
        # fichero JSON
        models_configs = read_json_config(json_config_path="./morphx/configs/models_configs.json")

        # CReamos el logger
        logger = create_logger(logger_file_name="PPT-GENERATOR")

        # Initialize the model for the slides
        slides_model = SlideModel(model_config=models_configs["slides_model"]["model_id"], logger=logger)

        # Initialize generator
        generator = PPTGenerator(slides_model=slides_model, logger=logger)
        
        # Validamos si el contenido cumple con la plantilla
        validate_slide_content(content=PPT_INPUT_CONTENT, logger=logger)

        # Generate presentation
        output_path = generator.generate_ppt(
            content_text=PPT_INPUT_CONTENT, 
            output_filename="my_presentation.pptx", 
            use_llm_enhancement=True
        )
        
        logger.info(f"SUCCESS! Presentation created: {output_path}")
        
    except PPTGeneratorError as e:
        logger.error(f"Generator error: {str(e)}")
        
    except Exception as e:
        logger.error(f"Unexpected error: {str(e)}")